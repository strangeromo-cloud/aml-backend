#!/usr/bin/env python3
"""Insert a Goals & Benefits slide (shapes+text only) as slide 2 of an existing deck,
preserving all the user's other edits."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x10, 0x2C, 0x3A)
INK = RGBColor(0x17, 0x31, 0x3F)
TEAL = RGBColor(0x1B, 0xA3, 0x9C)
MUTED = RGBColor(0x71, 0x83, 0x8F)
LINE = RGBColor(0xDC, 0xE8, 0xEE)
PANEL = RGBColor(0xF4, 0xFA, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _no_autofit(tf):
    # disable autosize so our font sizes stick
    pass


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Microsoft YaHei", line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    first = True
    for line in runs:  # runs = list of lines; each line = list of (text, size, color, bold)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for (text, size, color, bold) in line:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def rrect(slide, x, y, w, h, fill, line=None, line_w=1.0, radius=0.12):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    # adjust corner radius
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def build_goals_slide(prs, lang):
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    # white background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    font = "Microsoft YaHei" if lang == "cn" else "Calibri"
    W = 13.333

    # left teal accent bar (match cover)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.28), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False

    if lang == "cn":
        kick = "目标与收益"
        ttl = "让法务前置发现潜在洗钱行为"
        sub = "从「事后追查」转向「事前识别」—— 在付款出账前把风险拦下来。"
        cards = [
            ("1", "前置识别", "付款进行中即识别潜在洗钱，把风险拦在出款前，而非事后审计。"),
            ("2", "聚焦高价值", "海量付款里把最该看的顶到前面，法务有限产能用在刀刃上。"),
            ("3", "可解释 · 可审计", "每个告警都说得清「为什么」，分级与变更全程留痕，经得起监管追问。"),
            ("4", "越用越准 · 提效", "AI 调查取证提效；每次复核反哺系统，识别能力持续提升。"),
        ]
    else:
        kick = "GOALS & BENEFITS"
        ttl = "Help Legal catch potential laundering early"
        sub = "Shift from after-the-fact investigation to up-front detection — stop risk before the payment goes out."
        cards = [
            ("1", "Proactive detection", "Flag potential laundering while the payment is in flight — hold risk before payout, not audit it afterward."),
            ("2", "Focus on high value", "Surface the most-worth-reviewing few; spend Legal's limited capacity where it counts."),
            ("3", "Explainable & auditable", "Every alert says “why”; tiering and every change are fully logged — stands up to regulators."),
            ("4", "Sharper & faster", "AI speeds up evidence-gathering; every review feeds back and keeps improving detection."),
        ]

    add_text(slide, 0.6, 0.5, 10, 0.3, [[(kick, 12, TEAL, True)]], font=font)
    add_text(slide, 0.6, 0.82, 12.1, 0.9, [[(ttl, 30, INK, True)]], font=font)
    add_text(slide, 0.62, 1.78, 12.0, 0.6, [[(sub, 14, MUTED, False)]], font=font, line_spacing=1.15)

    cw, ch = 5.86, 1.95
    xs = [0.6, 6.88]
    ys = [2.6, 4.7]
    for i, (num, t, body) in enumerate(cards):
        x = xs[i % 2]
        y = ys[i // 2]
        rrect(slide, x, y, cw, ch, PANEL, LINE, 1.0, 0.10)
        # number badge
        bsp = rrect(slide, x + 0.3, y + 0.32, 0.62, 0.62, TEAL, None, radius=0.25)
        btf = bsp.text_frame; btf.word_wrap = False
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = num; br.font.size = Pt(20); br.font.bold = True
        br.font.color.rgb = WHITE; br.font.name = font
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # title + body
        add_text(slide, x + 1.12, y + 0.3, cw - 1.4, 0.5, [[(t, 16.5, INK, True)]], font=font, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + 0.34, y + 1.02, cw - 0.65, 0.8, [[(body, 12.5, MUTED, False)]], font=font, line_spacing=1.18)

    # move new slide to position 2 (index 1)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(1, ids[-1])


import sys
src, lang = sys.argv[1], sys.argv[2]
prs = Presentation(src)
build_goals_slide(prs, lang)
prs.save(src)
print(f"Inserted goals slide into {src} ({lang})")
