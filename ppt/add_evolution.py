#!/usr/bin/env python3
"""Insert a 'Two evolution engines' slide right after the self-evolution slide,
preserving all other (hand-edited) content. Shapes+text only."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x10, 0x2C, 0x3A); INK = RGBColor(0x17, 0x31, 0x3F)
TEAL = RGBColor(0x1B, 0xA3, 0x9C); AMBER = RGBColor(0xC9, 0x86, 0x1A)
MUTED = RGBColor(0x71, 0x83, 0x8F); LINE = RGBColor(0xDC, 0xE8, 0xEE)
PANEL = RGBColor(0xF4, 0xFA, 0xFB); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ARROWC = RGBColor(0x9F, 0xB6, 0xC0)


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei", ls=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04); tf.margin_top = tf.margin_bottom = Inches(0.02)
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align
        if ls: p.line_spacing = ls
        for (text, size, color, bold) in line:
            r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font
    return tb


def rrect(slide, x, y, w, h, fill, line=None, lw=1.0, radius=0.1):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    try: sp.adjustments[0] = radius
    except Exception: pass
    return sp


def box(slide, x, y, w, h, text, font):
    rrect(slide, x, y, w, h, PANEL, LINE, 1.0, 0.12)
    add_text(slide, x + 0.06, y, w - 0.12, h, [[(text, 12.5, INK, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=font, ls=1.0)


def arrow(slide, x, ycen, font):
    add_text(slide, x, ycen - 0.22, 0.42, 0.44, [[("→", 20, ARROWC, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=font)


def fbbar(slide, y, color, runs, font):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(12.13), Inches(0.62))
    sp.fill.solid(); sp.fill.fore_color.rgb = PANEL; sp.line.color.rgb = color; sp.line.width = Pt(1.2)
    sp.line.dash_style = None; sp.shadow.inherit = False
    try: sp.adjustments[0] = 0.18
    except Exception: pass
    # dashed line
    ln = sp.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
    add_text(slide, 0.95, y, 11.5, 0.62, [runs], anchor=MSO_ANCHOR.MIDDLE, font=font, ls=1.05)


def build(prs, lang):
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    for ph in list(slide.placeholders): ph._element.getparent().remove(ph._element)
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    font = "Microsoft YaHei" if lang == "cn" else "Calibri"
    # left teal accent
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.28), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False

    if lang == "cn":
        kick, ttl = "进化机制", "两套进化引擎"
        sub = "系统在两个方向上持续进化，但都受可审计闸门约束。"
        laneA = "进化 A · 调查越查越准 / 越全"
        boxesA = ["AI 调查", "法务复核", "沉淀调查 Playbook"]
        capA = [("回流：", 12.5, TEAL, True), ("每次复核的「漏了什么 / 查错了什么」沉淀成调查套路，下次查得更准、更全。", 12.5, INK, False)]
        laneB = "进化 B · 不断发现新 T1 / T2 规则"
        boxesB = ["Agent 挖候选信号", "PPV 统计校准", "人工审批", "新 T1 / T2 规则"]
        capB = [("闸门：", 12.5, AMBER, True), ("Agent 只提候选；达标 + 审批后才成为正式规则 —— 自我进化在闸门内。", 12.5, INK, False)]
    else:
        kick, ttl = "EVOLUTION", "Two evolution engines"
        sub = "The system evolves in two directions — both behind an auditable gate."
        laneA = "Loop A · Investigations get more accurate & complete"
        boxesA = ["AI investigation", "Legal review", "Update investigation playbook"]
        capA = [("Feedback: ", 12.5, TEAL, True), ("each review's misses become reusable playbooks — next time digs more accurately and completely.", 12.5, INK, False)]
        laneB = "Loop B · Keep discovering new T1 / T2 rules"
        boxesB = ["Agent mines candidates", "PPV calibration", "Human approval", "New T1 / T2 rule"]
        capB = [("Gate: ", 12.5, AMBER, True), ("the agent only proposes candidates; they go live only after they prove out + sign-off — evolution stays gated.", 12.5, INK, False)]

    add_text(slide, 0.6, 0.5, 10, 0.3, [[(kick, 12, TEAL, True)]], font=font)
    add_text(slide, 0.6, 0.82, 12.1, 0.9, [[(ttl, 30, INK, True)]], font=font)
    add_text(slide, 0.62, 1.75, 12.0, 0.5, [[(sub, 14, MUTED, False)]], font=font, ls=1.15)

    # Lane A — 3 boxes
    add_text(slide, 0.7, 2.35, 11, 0.32, [[(laneA, 13, TEAL, True)]], font=font)
    ay, ah = 2.72, 0.9
    aw, astep = 3.6, 4.05
    axs = [0.7, 0.7 + astep, 0.7 + 2 * astep]
    for i, b in enumerate(boxesA):
        box(slide, axs[i], ay, aw, ah, b, font)
        if i < len(boxesA) - 1: arrow(slide, axs[i] + aw + 0.02, ay + ah / 2, font)
    fbbar(slide, ay + ah + 0.12, TEAL, capA, font)

    # Lane B — 4 boxes
    by_lbl = 4.62
    add_text(slide, 0.75, by_lbl, 11, 0.32, [[(laneB, 13, AMBER, True)]], font=font)
    by, bh = 4.99, 0.9
    bw, bstep = 2.55, 2.95
    bxs = [0.75 + i * bstep for i in range(4)]
    for i, b in enumerate(boxesB):
        box(slide, bxs[i], by, bw, bh, b, font)
        if i < len(boxesB) - 1: arrow(slide, bxs[i] + bw + 0.0, by + bh / 2, font)
    fbbar(slide, by + bh + 0.12, AMBER, capB, font)

    # find self-evolution slide index, move new slide right after it
    def is_selfevo(sl):
        # match the self-evolution slide by its unique kicker (avoids the cover
        # subtitle and the summary, which also mention 自我进化 / Hermes).
        txt = " ".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)
        return ("重点创新" in txt) or ("SPOTLIGHT" in txt)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    new_id = ids[-1]
    slides = list(prs.slides)
    target = None
    for i, sl in enumerate(slides[:-1]):   # skip the just-added slide (last)
        if is_selfevo(sl):
            target = i
            break                          # first match = self-evolution, not the summary
    sldIdLst.remove(new_id)
    insert_at = (target + 1) if target is not None else len(ids) - 1
    sldIdLst.insert(insert_at, new_id)


src, lang = sys.argv[1], sys.argv[2]
prs = Presentation(src)
build(prs, lang)
prs.save(src)
print(f"Inserted evolution slide into {src} ({lang})")
